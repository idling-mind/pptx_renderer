import copy

from pptx.presentation import Presentation
from pptx.slide import Slide


def para_text_replace(para, find_string, replace_string):
    """Function to replace text in a paragraph

    This function replaces text in a paragraph while respecting the formatting.

    Args:
        para (pptx.shapes.paragraph.Paragraph): Paragraph to replace text in.
        find_string (str): String to find in the paragraph.
        replace_string (str): String to replace the find_string with.

    Returns:
        None
    """
    find_string = str(find_string)
    replace_string = str(replace_string)
    starting_pos = para.text.find(find_string)
    if starting_pos == -1:
        return  # text not in paragraph
    txt_prev = ""
    for run in para.runs:
        if len(txt_prev) <= starting_pos < len(txt_prev) + len(run.text):
            if run.text.find(find_string) != -1:  # text in run, replace
                run.text = run.text.replace(find_string, replace_string)
                return
            else:  # text no in "run"
                txt_prev = txt_prev + run.text
                run.text = run.text[: starting_pos - len(txt_prev)] + replace_string
        elif starting_pos < len(txt_prev) and starting_pos + len(find_string) >= len(
            txt_prev
        ) + len(run.text):
            txt_prev = txt_prev + run.text
            run.text = ""
        elif (
            len(txt_prev)
            < starting_pos + len(find_string)
            < len(txt_prev) + len(run.text)
        ):
            txt_prev = txt_prev + run.text
            run.text = run.text[starting_pos + len(find_string) - len(txt_prev) :]
        else:
            txt_prev += run.text


def fix_quotes(input_string: str) -> str:
    """Replace unicode quotes (inserted by powerpoint) with ascii quotes.

    Args:
        input_string (str): String to fix quotes in.
    
    Returns:
        str: String with fixed quotes.
    """
    return (
        input_string.replace("’", "'")
        .replace("‘", "'")
        .replace("“", '"')
        .replace("”", '"')
    )

def copy_slide(source_ppt: Presentation, target_ppt: Presentation, slide: Slide) -> Slide:
    """Duplicate each slide in prs2 and "moves" it into prs1.
    Adds slides to the end of the presentation

    Args:
        source_ppt (Presentation): Source presentation.
        target_ppt (Presentation): Target presentation.
        slide (Slide): Slide to copy.

    Returns:
        Slide: Slide that was copied
    """
    layout = source_ppt.slide_layouts.index(slide.slide_layout)
    new_slide = target_ppt.slides.add_slide(target_ppt.slide_layouts[layout])
    for new_ph, old_ph in zip(new_slide.placeholders, slide.placeholders):
        new_ph.text = old_ph.text
    for shape in slide.shapes:
        if shape.is_placeholder:
            continue
        newel = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
    return new_slide

def clear_presentation(prs: Presentation):
    """Clears all slides from a presentation"""
    for i in range(len(prs.slides)-1, -1, -1): 
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]